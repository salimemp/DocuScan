"""Document export generators.

These pure functions take a document dict (as stored in MongoDB) and return
bytes for one of the 12 supported export formats: PDF, DOCX, PPTX, TXT,
HTML, JSON, Markdown, PNG/JPEG/TIFF/BMP/WebP (via generate_image_export),
XLSX, EPUB, MOBI, SVG.

They were extracted from the monolithic `server.py` so that `server.py`
can shrink and so that we can independently unit-test each generator.

The `safe_latin` helper used by several generators also lives here.
"""
from __future__ import annotations

import io
import json
import uuid
from datetime import datetime, timezone


def safe_latin(text: str) -> str:
    """Encode text for libraries (like fpdf) that only support latin-1.

    Characters outside latin-1 are replaced with `?`. This is the same
    fallback used in the original server.py — fine for short titles and
    metadata; the full OCR text should already be UTF-8-safe by the time
    it reaches us.
    """
    return str(text or "").encode("latin-1", errors="replace").decode("latin-1")


# ── Format → generator dispatch table ───────────────────────────────────────
# Used by the export endpoint in server.py (see generate_export below).
GENERATORS = {}


def _register(*formats: str):
    """Decorator that registers a generator under one or more format names."""
    def wrap(fn):
        for fmt in formats:
            GENERATORS[fmt.lower()] = fn
        return fn
    return wrap


@_register("pdf")
def generate_pdf(doc: dict) -> bytes:
    from fpdf import FPDF
    from fpdf.enums import XPos, YPos
    TYPE_LABELS = {
        "passport": "PASSPORT", "national_id": "NATIONAL ID", "drivers_license": "DRIVER'S LICENSE",
        "invoice": "INVOICE", "receipt": "RECEIPT", "business_card": "BUSINESS CARD",
        "contract": "CONTRACT", "bank_statement": "BANK STATEMENT", "medical_record": "MEDICAL RECORD",
        "prescription": "PRESCRIPTION", "handwritten_note": "HANDWRITTEN NOTE", "certificate": "CERTIFICATE",
        "legal_document": "LEGAL DOCUMENT", "academic_transcript": "ACADEMIC TRANSCRIPT",
        "tax_document": "TAX DOCUMENT", "insurance_document": "INSURANCE DOCUMENT",
        "utility_bill": "UTILITY BILL", "general_document": "DOCUMENT",
    }
    type_label = TYPE_LABELS.get(doc.get("document_type", ""), "DOCUMENT")
    pdf = FPDF(orientation="P", unit="mm", format="A4")
    pdf.set_auto_page_break(auto=True, margin=18)
    pdf.add_page()
    pdf.set_fill_color(37, 99, 235)
    pdf.rect(0, 0, 210, 22, "F")
    pdf.set_font("Helvetica", "B", 12)
    pdf.set_text_color(255, 255, 255)
    pdf.set_xy(10, 6)
    pdf.cell(100, 10, "DocScan Pro", new_x=XPos.RIGHT, new_y=YPos.TOP)
    pdf.set_x(110)
    pdf.cell(0, 10, safe_latin(type_label), new_x=XPos.LMARGIN, new_y=YPos.NEXT, align="R")
    pdf.set_text_color(30, 30, 30)
    pdf.set_xy(15, 30)
    pdf.set_font("Helvetica", "B", 18)
    pdf.multi_cell(180, 10, safe_latin(doc.get("title", "Untitled Document")))
    pdf.set_font("Helvetica", "", 9)
    pdf.set_text_color(100, 100, 100)
    pdf.ln(2)
    lang = safe_latin(doc.get("detected_language", "Unknown"))
    conf = f"{int(doc.get('confidence', 0) * 100)}%"
    pdf.cell(65, 6, f"Language: {lang}")
    pdf.cell(0, 6, f"Confidence: {conf}", new_x=XPos.LMARGIN, new_y=YPos.NEXT)
    y = pdf.get_y() + 3
    pdf.set_draw_color(37, 99, 235)
    pdf.line(15, y, 195, y)
    pdf.ln(5)
    if doc.get("summary"):
        pdf.set_font("Helvetica", "I", 10)
        pdf.set_text_color(70, 70, 70)
        pdf.multi_cell(180, 6, safe_latin(doc["summary"]))
        pdf.ln(4)
    if doc.get("formatted_output"):
        pdf.set_font("Helvetica", "B", 12)
        pdf.set_text_color(37, 99, 235)
        pdf.cell(0, 8, "Extracted Content", new_x=XPos.LMARGIN, new_y=YPos.NEXT)
        pdf.ln(3)
        pdf.set_font("Courier", "", 8)
        pdf.set_text_color(40, 40, 40)
        for line in doc["formatted_output"].split("\n")[:100]:
            if line.strip():
                pdf.multi_cell(180, 5, safe_latin(line))
            else:
                pdf.ln(2)
    return bytes(pdf.output())


@_register("txt")
def generate_txt(doc: dict) -> bytes:
    lines = ["DocScan Pro — Document Export", "=" * 60,
             f"Title:      {doc.get('title', '')}",
             f"Type:       {doc.get('document_type', '').replace('_', ' ').title()}",
             f"Language:   {doc.get('detected_language', '')}",
             f"Confidence: {int(doc.get('confidence', 0) * 100)}%", ""]
    if doc.get("summary"):
        lines += ["SUMMARY", "-" * 40, doc["summary"], ""]
    if doc.get("formatted_output"):
        lines += ["CONTENT", "-" * 40, doc["formatted_output"], ""]
    return "\n".join(lines).encode("utf-8")


@_register("docx")
def generate_docx(doc: dict) -> bytes:
    from docx import Document as DocxDoc
    d = DocxDoc()
    d.add_heading(doc.get("title", "Untitled"), 0)
    if doc.get("summary"):
        p = d.add_paragraph(doc["summary"])
        p.runs[0].italic = True
    if doc.get("formatted_output"):
        d.add_heading("Content", 1)
        for line in doc["formatted_output"].split("\n"):
            if line.strip():
                d.add_paragraph(line)
    buf = io.BytesIO()
    d.save(buf)
    return buf.getvalue()


@_register("pptx")
def generate_pptx(doc: dict) -> bytes:
    from pptx import Presentation
    prs = Presentation()
    sl = prs.slides.add_slide(prs.slide_layouts[0])
    sl.shapes.title.text = doc.get("title", "Document")
    sl.placeholders[1].text = doc.get("summary", "")
    if doc.get("formatted_output"):
        sl2 = prs.slides.add_slide(prs.slide_layouts[1])
        sl2.shapes.title.text = "Content"
        sl2.placeholders[1].text = doc["formatted_output"][:2000]
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


@_register("png", "jpg", "jpeg", "tiff", "bmp", "webp")
def generate_image_export(doc: dict, fmt: str = "JPEG") -> bytes:
    from PIL import Image, ImageDraw
    W, H = 900, 1200
    img = Image.new("RGB", (W, H), color=(255, 255, 255))
    draw = ImageDraw.Draw(img)
    draw.rectangle([0, 0, W, 70], fill=(37, 99, 235))
    draw.text((W // 2, 35), "DocScan Pro", fill=(255, 255, 255), anchor="mm")
    draw.text((W // 2, 100), (doc.get("title", "") or "")[:70], fill=(20, 20, 20), anchor="mt")
    y = 150
    if doc.get("formatted_output"):
        for line in doc["formatted_output"].split("\n")[:30]:
            draw.text((40, y), line[:80], fill=(60, 60, 60))
            y += 20
            if y > H - 100:
                break
    buf = io.BytesIO()
    fmt_upper = fmt.upper()
    if fmt_upper == "JPG":
        fmt_upper = "JPEG"
    img.save(buf, format=fmt_upper)
    return buf.getvalue()


@_register("xlsx")
def generate_xlsx(doc: dict) -> bytes:
    import xlsxwriter
    buf = io.BytesIO()
    wb = xlsxwriter.Workbook(buf, {"in_memory": True})

    title_fmt = wb.add_format({"bold": True, "font_size": 16, "font_color": "#2563EB"})
    header_fmt = wb.add_format({"bold": True, "bg_color": "#2563EB", "font_color": "white", "border": 1})
    cell_fmt = wb.add_format({"text_wrap": True, "valign": "top", "border": 1})

    ws_info = wb.add_worksheet("Document Info")
    ws_info.set_column("A:A", 20)
    ws_info.set_column("B:B", 60)

    ws_info.write("A1", doc.get("title", "Untitled"), title_fmt)
    ws_info.write("A3", "Property", header_fmt)
    ws_info.write("B3", "Value", header_fmt)

    info_rows = [
        ("Document Type", doc.get("document_type", "").replace("_", " ").title()),
        ("Language", doc.get("detected_language", "Unknown")),
        ("Confidence", f"{int(doc.get('confidence', 0) * 100)}%"),
        ("Pages", str(doc.get("pages_count", 1))),
        ("Created", doc.get("created_at", "")),
    ]
    for i, (prop, val) in enumerate(info_rows, start=4):
        ws_info.write(f"A{i}", prop, cell_fmt)
        ws_info.write(f"B{i}", val, cell_fmt)

    if doc.get("formatted_output"):
        ws_content = wb.add_worksheet("Content")
        ws_content.set_column("A:A", 100)
        ws_content.write("A1", "Extracted Content", title_fmt)
        row = 2
        for line in doc["formatted_output"].split("\n"):
            ws_content.write(f"A{row}", line, cell_fmt)
            row += 1

    if doc.get("summary"):
        ws_summary = wb.add_worksheet("Summary")
        ws_summary.set_column("A:A", 80)
        ws_summary.write("A1", "Document Summary", title_fmt)
        ws_summary.write("A3", doc["summary"], cell_fmt)

    wb.close()
    return buf.getvalue()


@_register("html")
def generate_html(doc: dict) -> bytes:
    html = f'''<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{doc.get('title', 'Document')}</title>
    <style>
        body {{ font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, sans-serif; max-width: 800px; margin: 0 auto; padding: 20px; background: #f8fafc; }}
        .header {{ background: #2563eb; color: white; padding: 20px; border-radius: 12px; margin-bottom: 20px; }}
        .header h1 {{ margin: 0 0 10px 0; }}
        .meta {{ display: flex; gap: 20px; font-size: 14px; opacity: 0.9; }}
        .section {{ background: white; padding: 20px; border-radius: 12px; margin-bottom: 16px; box-shadow: 0 1px 3px rgba(0,0,0,0.1); }}
        .section h2 {{ color: #2563eb; margin-top: 0; border-bottom: 2px solid #e5e7eb; padding-bottom: 10px; }}
        .content {{ white-space: pre-wrap; font-family: 'SF Mono', Monaco, monospace; font-size: 13px; line-height: 1.6; }}
        .tag {{ display: inline-block; background: #dbeafe; color: #2563eb; padding: 4px 12px; border-radius: 16px; font-size: 12px; margin: 4px; }}
        .footer {{ text-align: center; color: #6b7280; font-size: 12px; margin-top: 20px; }}
    </style>
</head>
<body>
    <div class="header">
        <h1>{doc.get('title', 'Untitled Document')}</h1>
        <div class="meta">
            <span>📄 {doc.get('document_type', 'document').replace('_', ' ').title()}</span>
            <span>🌐 {doc.get('detected_language', 'Unknown')}</span>
            <span>✓ {int(doc.get('confidence', 0) * 100)}% confidence</span>
        </div>
    </div>
'''
    if doc.get("summary"):
        html += f'''    <div class="section">
        <h2>📋 Summary</h2>
        <p>{doc['summary']}</p>
    </div>
'''
    if doc.get("formatted_output"):
        html += f'''    <div class="section">
        <h2>📝 Content</h2>
        <div class="content">{doc['formatted_output']}</div>
    </div>
'''
    if doc.get("tags"):
        tags_html = "".join([f'<span class="tag">#{tag}</span>' for tag in doc["tags"]])
        html += f'''    <div class="section">
        <h2>🏷️ Tags</h2>
        {tags_html}
    </div>
'''
    html += '''    <div class="footer">
        <p>Exported from DocScan Pro</p>
    </div>
</body>
</html>'''
    return html.encode("utf-8")


@_register("json")
def generate_json_export(doc: dict) -> bytes:
    created_at = doc.get("created_at", "")
    if hasattr(created_at, "isoformat"):
        created_at = created_at.isoformat()
    elif created_at and not isinstance(created_at, str):
        created_at = str(created_at)

    export_data = {
        "document": {
            "title": doc.get("title", ""),
            "type": doc.get("document_type", ""),
            "language": doc.get("detected_language", ""),
            "confidence": doc.get("confidence", 0),
            "pages_count": doc.get("pages_count", 1),
            "created_at": created_at,
        },
        "content": {
            "summary": doc.get("summary", ""),
            "formatted_output": doc.get("formatted_output", ""),
            "raw_text": doc.get("raw_text", ""),
        },
        "metadata": {
            "tags": doc.get("tags", []),
            "is_locked": doc.get("is_locked", False),
            "comments_count": len(doc.get("comments", [])),
            "signatures_count": len(doc.get("signatures", [])),
        },
        "export_info": {
            "exported_at": datetime.now(timezone.utc).isoformat(),
            "source": "DocScan Pro",
            "version": "5.0",
        },
    }
    return json.dumps(export_data, indent=2, ensure_ascii=False).encode("utf-8")


@_register("md", "markdown")
def generate_markdown(doc: dict) -> bytes:
    md = f'''# {doc.get('title', 'Untitled Document')}

---

**Document Type:** {doc.get('document_type', 'document').replace('_', ' ').title()}  
**Language:** {doc.get('detected_language', 'Unknown')}  
**Confidence:** {int(doc.get('confidence', 0) * 100)}%  
**Pages:** {doc.get('pages_count', 1)}

---

'''
    if doc.get("summary"):
        md += f'''## 📋 Summary

{doc['summary']}

'''
    if doc.get("formatted_output"):
        md += f'''## 📝 Content

```
{doc['formatted_output']}
```

'''
    if doc.get("tags"):
        tags = " ".join([f"`#{tag}`" for tag in doc["tags"]])
        md += f'''## 🏷️ Tags

{tags}

'''
    md += '''---

*Exported from DocScan Pro*
'''
    return md.encode("utf-8")


@_register("epub")
def generate_epub(doc: dict) -> bytes:
    from ebooklib import epub
    book = epub.EpubBook()
    book.set_identifier(doc.get("id", str(uuid.uuid4())))
    book.set_title(doc.get("title", "Document"))
    book.set_language(doc.get("detected_language", "en")[:2].lower() or "en")
    book.add_author("DocScan Pro")

    c1 = epub.EpubHtml(title=doc.get("title", "Document"), file_name="content.xhtml", lang="en")
    content_html = f'''<html>
<head><title>{doc.get('title', 'Document')}</title></head>
<body>
<h1>{doc.get('title', 'Document')}</h1>
<p><em>Type: {doc.get('document_type', '').replace('_', ' ').title()}</em></p>
'''
    if doc.get("summary"):
        content_html += f'<h2>Summary</h2><p>{doc["summary"]}</p>'
    if doc.get("formatted_output"):
        content_html += f'<h2>Content</h2><pre>{doc["formatted_output"]}</pre>'
    content_html += "</body></html>"
    c1.content = content_html

    book.add_item(c1)
    book.spine = ["nav", c1]
    book.add_item(epub.EpubNcx())
    book.add_item(epub.EpubNav())

    buf = io.BytesIO()
    epub.write_epub(buf, book)
    return buf.getvalue()


@_register("mobi")
def generate_mobi(doc: dict) -> bytes:
    """MOBI export — fall back to EPUB since true MOBI needs Amazon's Kindlegen."""
    return generate_epub(doc)


@_register("svg")
def generate_svg(doc: dict) -> bytes:
    title = doc.get("title", "Document")[:50]
    doc_type = doc.get("document_type", "document").replace("_", " ").title()
    content_lines = (doc.get("formatted_output", "") or "").split("\n")[:20]

    svg = f'''<?xml version="1.0" encoding="UTF-8"?>
<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 600 800" width="600" height="800">
  <defs>
    <linearGradient id="headerGrad" x1="0%" y1="0%" x2="100%" y2="0%">
      <stop offset="0%" style="stop-color:#2563eb"/>
      <stop offset="100%" style="stop-color:#3b82f6"/>
    </linearGradient>
  </defs>

  <!-- Background -->
  <rect width="600" height="800" fill="#f8fafc"/>

  <!-- Header -->
  <rect width="600" height="80" fill="url(#headerGrad)"/>
  <text x="300" y="35" text-anchor="middle" fill="white" font-family="Arial" font-size="18" font-weight="bold">DocScan Pro</text>
  <text x="300" y="60" text-anchor="middle" fill="rgba(255,255,255,0.8)" font-family="Arial" font-size="12">{doc_type}</text>

  <!-- Title -->
  <text x="30" y="120" fill="#1e3a5f" font-family="Arial" font-size="20" font-weight="bold">{title}</text>

  <!-- Content -->
  <rect x="20" y="140" width="560" height="620" fill="white" stroke="#e5e7eb" rx="8"/>
'''
    y = 170
    for line in content_lines:
        escaped = line[:70].replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
        svg += f'  <text x="35" y="{y}" fill="#374151" font-family="Monaco, monospace" font-size="11">{escaped}</text>\n'
        y += 22
        if y > 720:
            break

    svg += '''
  <!-- Footer -->
  <text x="300" y="780" text-anchor="middle" fill="#9ca3af" font-family="Arial" font-size="10">Exported from DocScan Pro</text>
</svg>'''
    return svg.encode("utf-8")


def generate_export(doc: dict, fmt: str) -> tuple[bytes, str]:
    """Dispatch helper for the export endpoint.

    Returns (bytes, content_type) so the endpoint doesn't need to know
    about PIL/format mapping. Raises KeyError on unknown formats — the
    caller (server.py) catches that and returns 400.
    """
    fmt_lc = fmt.lower()
    if fmt_lc not in GENERATORS:
        raise KeyError(fmt_lc)
    content = GENERATORS[fmt_lc](doc)
    content_type = {
        "pdf": "application/pdf",
        "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "html": "text/html; charset=utf-8",
        "json": "application/json",
        "txt": "text/plain; charset=utf-8",
        "md": "text/markdown; charset=utf-8",
        "markdown": "text/markdown; charset=utf-8",
        "epub": "application/epub+zip",
        "mobi": "application/x-mobipocket-ebook",
        "svg": "image/svg+xml",
        "png": "image/png",
        "jpg": "image/jpeg",
        "jpeg": "image/jpeg",
        "tiff": "image/tiff",
        "bmp": "image/bmp",
        "webp": "image/webp",
    }.get(fmt_lc, "application/octet-stream")
    return content, content_type