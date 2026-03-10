from fastapi import FastAPI, APIRouter, HTTPException, Query, BackgroundTasks, UploadFile, File
from fastapi.responses import JSONResponse
from dotenv import load_dotenv
from starlette.middleware.cors import CORSMiddleware
from motor.motor_asyncio import AsyncIOMotorClient
import os, logging, json, re, uuid, base64, io, hashlib, cv2, numpy as np
from pathlib import Path
from pydantic import BaseModel, Field
from typing import List, Optional, Any, Dict
from datetime import datetime, timezone
from PIL import Image, ImageDraw, ImageFont, ImageFilter, ImageEnhance
from emergentintegrations.llm.chat import LlmChat, UserMessage, ImageContent
import resend

# Import auth and subscription modules
from auth import auth_router
from subscriptions import subscription_router
from document_security import security_router

ROOT_DIR = Path(__file__).parent
load_dotenv(ROOT_DIR / '.env')

mongo_url = os.environ['MONGO_URL']
client = AsyncIOMotorClient(mongo_url)
db = client[os.environ['DB_NAME']]

# Initialize Resend
resend.api_key = os.environ.get('RESEND_API_KEY', '')

app = FastAPI()
api_router = APIRouter(prefix="/api")

logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

# ── Helpers ────────────────────────────────────────────────────────────────
def get_api_key() -> str:
    key = os.environ.get('GEMINI_API_KEY') or os.environ.get('EMERGENT_LLM_KEY')
    if not key:
        raise HTTPException(500, "No AI API key configured")
    return key

def strip_b64_prefix(s: str) -> str:
    return s.split(',', 1)[1] if ',' in s and s.startswith('data:') else s

def safe_latin(text: str) -> str:
    return str(text or '').encode('latin-1', errors='replace').decode('latin-1')

def hash_password(password: str) -> str:
    return hashlib.sha256(password.encode()).hexdigest()

def verify_password(password: str, hashed: str) -> bool:
    return hash_password(password) == hashed

def b64_to_cv2(b64_string: str) -> np.ndarray:
    """Convert base64 to OpenCV image"""
    img_data = base64.b64decode(strip_b64_prefix(b64_string))
    nparr = np.frombuffer(img_data, np.uint8)
    return cv2.imdecode(nparr, cv2.IMREAD_COLOR)

def cv2_to_b64(img: np.ndarray, fmt: str = 'JPEG') -> str:
    """Convert OpenCV image to base64"""
    _, buffer = cv2.imencode(f'.{fmt.lower()}', img)
    return base64.b64encode(buffer).decode()

def pil_to_b64(img: Image.Image, fmt: str = 'JPEG') -> str:
    """Convert PIL image to base64"""
    buf = io.BytesIO()
    img.save(buf, format=fmt, quality=90)
    return base64.b64encode(buf.getvalue()).decode()

def b64_to_pil(b64_string: str) -> Image.Image:
    """Convert base64 to PIL image"""
    img_data = base64.b64decode(strip_b64_prefix(b64_string))
    return Image.open(io.BytesIO(img_data))

# ── Email Templates ────────────────────────────────────────────────────────
def get_email_template(template_type: str, data: dict) -> tuple:
    """Generate professional HTML email templates"""
    
    base_style = """
    <style>
        body { font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, sans-serif; margin: 0; padding: 0; background: #F3F4F6; }
        .container { max-width: 600px; margin: 0 auto; background: #FFFFFF; border-radius: 16px; overflow: hidden; box-shadow: 0 4px 6px rgba(0,0,0,0.1); }
        .header { background: linear-gradient(135deg, #2563EB 0%, #1D4ED8 100%); padding: 32px; text-align: center; }
        .header h1 { color: #FFFFFF; font-size: 24px; margin: 0; }
        .header p { color: rgba(255,255,255,0.9); font-size: 14px; margin: 8px 0 0 0; }
        .content { padding: 32px; }
        .content h2 { color: #1F2937; font-size: 20px; margin: 0 0 16px 0; }
        .content p { color: #4B5563; font-size: 15px; line-height: 1.6; margin: 0 0 16px 0; }
        .btn { display: inline-block; background: #2563EB; color: #FFFFFF; text-decoration: none; padding: 14px 28px; border-radius: 10px; font-weight: 600; font-size: 15px; }
        .btn:hover { background: #1D4ED8; }
        .info-box { background: #F0F9FF; border-left: 4px solid #2563EB; padding: 16px; border-radius: 0 8px 8px 0; margin: 20px 0; }
        .info-box p { margin: 0; color: #1E40AF; }
        .footer { background: #F9FAFB; padding: 24px 32px; text-align: center; border-top: 1px solid #E5E7EB; }
        .footer p { color: #9CA3AF; font-size: 12px; margin: 0; }
        .document-card { background: #F9FAFB; border-radius: 12px; padding: 20px; margin: 20px 0; }
        .document-card h3 { color: #1F2937; font-size: 16px; margin: 0 0 8px 0; }
        .document-card span { color: #6B7280; font-size: 13px; }
    </style>
    """
    
    if template_type == 'signature_request':
        subject = f"✍️ Signature Request: {data.get('document_title', 'Document')}"
        html = f"""
        <!DOCTYPE html>
        <html>
        <head>{base_style}</head>
        <body>
            <div style="padding: 20px;">
                <div class="container">
                    <div class="header">
                        <h1>📝 DocScan Pro</h1>
                        <p>Signature Request</p>
                    </div>
                    <div class="content">
                        <h2>Hello {data.get('signer_name', 'there')}!</h2>
                        <p><strong>{data.get('requester_name', 'Someone')}</strong> has requested your signature on a document.</p>
                        
                        <div class="document-card">
                            <h3>📄 {data.get('document_title', 'Untitled Document')}</h3>
                            <span>Requested on {datetime.now().strftime('%B %d, %Y')}</span>
                        </div>
                        
                        {f'<div class="info-box"><p><strong>Message:</strong> {data.get("message")}</p></div>' if data.get('message') else ''}
                        
                        <p>Please review and sign the document at your earliest convenience.</p>
                        
                        <p style="text-align: center; margin-top: 24px;">
                            <a href="{data.get('action_url', '#')}" class="btn">Review & Sign Document</a>
                        </p>
                    </div>
                    <div class="footer">
                        <p>This email was sent by DocScan Pro on behalf of {data.get('requester_email', 'a user')}.</p>
                        <p style="margin-top: 8px;">© {datetime.now().year} DocScan Pro. All rights reserved.</p>
                    </div>
                </div>
            </div>
        </body>
        </html>
        """
        return subject, html
    
    elif template_type == 'comment_request':
        subject = f"💬 Review Request: {data.get('document_title', 'Document')}"
        html = f"""
        <!DOCTYPE html>
        <html>
        <head>{base_style}</head>
        <body>
            <div style="padding: 20px;">
                <div class="container">
                    <div class="header" style="background: linear-gradient(135deg, #059669 0%, #047857 100%);">
                        <h1>📝 DocScan Pro</h1>
                        <p>Review Request</p>
                    </div>
                    <div class="content">
                        <h2>Hello {data.get('reviewer_name', 'there')}!</h2>
                        <p><strong>{data.get('requester_name', 'Someone')}</strong> has requested your review and comments on a document.</p>
                        
                        <div class="document-card">
                            <h3>📄 {data.get('document_title', 'Untitled Document')}</h3>
                            <span>Review requested on {datetime.now().strftime('%B %d, %Y')}</span>
                        </div>
                        
                        {f'<div class="info-box" style="border-color: #059669;"><p style="color: #047857;"><strong>Message:</strong> {data.get("message")}</p></div>' if data.get('message') else ''}
                        
                        <p>Please review the document and add your comments.</p>
                        
                        <p style="text-align: center; margin-top: 24px;">
                            <a href="{data.get('action_url', '#')}" class="btn" style="background: #059669;">Review Document</a>
                        </p>
                    </div>
                    <div class="footer">
                        <p>This email was sent by DocScan Pro on behalf of {data.get('requester_email', 'a user')}.</p>
                        <p style="margin-top: 8px;">© {datetime.now().year} DocScan Pro. All rights reserved.</p>
                    </div>
                </div>
            </div>
        </body>
        </html>
        """
        return subject, html
    
    elif template_type == 'document_shared':
        subject = f"📎 Document Shared: {data.get('document_title', 'Document')}"
        html = f"""
        <!DOCTYPE html>
        <html>
        <head>{base_style}</head>
        <body>
            <div style="padding: 20px;">
                <div class="container">
                    <div class="header" style="background: linear-gradient(135deg, #7C3AED 0%, #6D28D9 100%);">
                        <h1>📝 DocScan Pro</h1>
                        <p>Document Shared With You</p>
                    </div>
                    <div class="content">
                        <h2>Hello {data.get('recipient_name', 'there')}!</h2>
                        <p><strong>{data.get('sender_name', 'Someone')}</strong> has shared a document with you.</p>
                        
                        <div class="document-card">
                            <h3>📄 {data.get('document_title', 'Untitled Document')}</h3>
                            <span>Shared on {datetime.now().strftime('%B %d, %Y at %I:%M %p')}</span>
                        </div>
                        
                        {f'<div class="info-box" style="border-color: #7C3AED;"><p style="color: #6D28D9;"><strong>Note:</strong> {data.get("message")}</p></div>' if data.get('message') else ''}
                        
                        <p style="text-align: center; margin-top: 24px;">
                            <a href="{data.get('action_url', '#')}" class="btn" style="background: #7C3AED;">View Document</a>
                        </p>
                    </div>
                    <div class="footer">
                        <p>This email was sent by DocScan Pro on behalf of {data.get('sender_email', 'a user')}.</p>
                        <p style="margin-top: 8px;">© {datetime.now().year} DocScan Pro. All rights reserved.</p>
                    </div>
                </div>
            </div>
        </body>
        </html>
        """
        return subject, html
    
    return "DocScan Pro Notification", "<p>You have a new notification from DocScan Pro.</p>"

async def send_email(to_email: str, template_type: str, data: dict):
    """Send email using Resend"""
    try:
        subject, html = get_email_template(template_type, data)
        
        params = {
            "from": "DocScan Pro <onboarding@resend.dev>",
            "to": [to_email],
            "subject": subject,
            "html": html,
        }
        
        email_response = resend.Emails.send(params)
        logger.info(f"Email sent successfully to {to_email}: {email_response}")
        return True
    except Exception as e:
        logger.error(f"Failed to send email to {to_email}: {e}")
        return False

# ── Image Processing Functions ────────────────────────────────────────────
def straighten_document(img: np.ndarray) -> np.ndarray:
    """Auto-straighten/deskew a document image"""
    gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
    gray = cv2.bitwise_not(gray)
    
    # Find coordinates of non-zero pixels
    coords = np.column_stack(np.where(gray > 0))
    
    if len(coords) < 10:
        return img
    
    # Get the minimum area rectangle
    angle = cv2.minAreaRect(coords)[-1]
    
    if angle < -45:
        angle = -(90 + angle)
    else:
        angle = -angle
    
    # Limit rotation to reasonable angles
    if abs(angle) > 15:
        angle = 0
    
    (h, w) = img.shape[:2]
    center = (w // 2, h // 2)
    M = cv2.getRotationMatrix2D(center, angle, 1.0)
    rotated = cv2.warpAffine(img, M, (w, h), flags=cv2.INTER_CUBIC, borderMode=cv2.BORDER_REPLICATE)
    
    return rotated

def enhance_document(img: np.ndarray) -> np.ndarray:
    """Enhance document image for better readability"""
    # Convert to LAB color space
    lab = cv2.cvtColor(img, cv2.COLOR_BGR2LAB)
    l, a, b = cv2.split(lab)
    
    # Apply CLAHE to L channel
    clahe = cv2.createCLAHE(clipLimit=2.0, tileGridSize=(8, 8))
    cl = clahe.apply(l)
    
    # Merge channels
    enhanced_lab = cv2.merge((cl, a, b))
    enhanced = cv2.cvtColor(enhanced_lab, cv2.COLOR_LAB2BGR)
    
    # Sharpen
    kernel = np.array([[-1, -1, -1], [-1, 9, -1], [-1, -1, -1]])
    enhanced = cv2.filter2D(enhanced, -1, kernel)
    
    return enhanced

def detect_document_edges(img: np.ndarray) -> list:
    """Detect document edges and return corner points"""
    gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
    blur = cv2.GaussianBlur(gray, (5, 5), 0)
    edges = cv2.Canny(blur, 50, 150)
    
    contours, _ = cv2.findContours(edges, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
    
    if not contours:
        return []
    
    # Find the largest contour
    largest = max(contours, key=cv2.contourArea)
    
    # Approximate the contour to a polygon
    epsilon = 0.02 * cv2.arcLength(largest, True)
    approx = cv2.approxPolyDP(largest, epsilon, True)
    
    if len(approx) == 4:
        return approx.reshape(4, 2).tolist()
    
    return []

def apply_perspective_transform(img: np.ndarray, points: list) -> np.ndarray:
    """Apply perspective transform to flatten document"""
    if len(points) != 4:
        return img
    
    pts = np.array(points, dtype=np.float32)
    
    # Order points: top-left, top-right, bottom-right, bottom-left
    rect = np.zeros((4, 2), dtype=np.float32)
    s = pts.sum(axis=1)
    rect[0] = pts[np.argmin(s)]
    rect[2] = pts[np.argmax(s)]
    diff = np.diff(pts, axis=1)
    rect[1] = pts[np.argmin(diff)]
    rect[3] = pts[np.argmax(diff)]
    
    # Calculate dimensions
    width_a = np.linalg.norm(rect[2] - rect[3])
    width_b = np.linalg.norm(rect[1] - rect[0])
    max_width = max(int(width_a), int(width_b))
    
    height_a = np.linalg.norm(rect[1] - rect[2])
    height_b = np.linalg.norm(rect[0] - rect[3])
    max_height = max(int(height_a), int(height_b))
    
    dst = np.array([
        [0, 0],
        [max_width - 1, 0],
        [max_width - 1, max_height - 1],
        [0, max_height - 1]
    ], dtype=np.float32)
    
    M = cv2.getPerspectiveTransform(rect, dst)
    warped = cv2.warpPerspective(img, M, (max_width, max_height))
    
    return warped

def add_watermark(img: Image.Image, text: str, opacity: float = 0.3) -> Image.Image:
    """Add text watermark to image"""
    watermark = Image.new('RGBA', img.size, (0, 0, 0, 0))
    draw = ImageDraw.Draw(watermark)
    
    # Calculate font size based on image size
    font_size = min(img.size) // 10
    try:
        font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", font_size)
    except:
        font = ImageFont.load_default()
    
    # Get text size
    bbox = draw.textbbox((0, 0), text, font=font)
    text_width = bbox[2] - bbox[0]
    text_height = bbox[3] - bbox[1]
    
    # Position diagonally across image
    x = (img.size[0] - text_width) // 2
    y = (img.size[1] - text_height) // 2
    
    # Draw text with opacity
    alpha = int(255 * opacity)
    draw.text((x, y), text, font=font, fill=(128, 128, 128, alpha))
    
    # Rotate watermark
    watermark = watermark.rotate(45, expand=False, center=(img.size[0]//2, img.size[1]//2))
    
    # Composite
    if img.mode != 'RGBA':
        img = img.convert('RGBA')
    
    return Image.alpha_composite(img, watermark)

def add_blur_region(img: Image.Image, x: int, y: int, width: int, height: int, intensity: int = 20) -> Image.Image:
    """Blur a specific region of the image"""
    # Crop the region
    region = img.crop((x, y, x + width, y + height))
    
    # Apply blur
    blurred = region.filter(ImageFilter.GaussianBlur(radius=intensity))
    
    # Paste back
    result = img.copy()
    result.paste(blurred, (x, y))
    
    return result

# ── Models ────────────────────────────────────────────────────────────────
class ScanRequest(BaseModel):
    images: List[str]
    mime_type: str = "image/jpeg"

class ImageProcessRequest(BaseModel):
    image: str  # base64
    operations: List[str] = []  # ['straighten', 'enhance', 'detect_edges']

class WatermarkRequest(BaseModel):
    image: str
    text: str
    opacity: float = 0.3

class BlurRequest(BaseModel):
    image: str
    x: int
    y: int
    width: int
    height: int
    intensity: int = 20

class AIAssistantRequest(BaseModel):
    document_id: Optional[str] = None
    message: str
    context: Optional[str] = None

class CloudSyncRequest(BaseModel):
    provider: str  # google_drive, dropbox, onedrive, box, icloud
    action: str  # upload, download, list
    document_id: Optional[str] = None
    folder_path: Optional[str] = None

class MeasurementRequest(BaseModel):
    image: str
    mode: str  # count, area
    points: Optional[List[Dict[str, float]]] = None

class Comment(BaseModel):
    id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    author: str = "Anonymous"
    author_email: Optional[str] = None
    content: str
    created_at: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))
    resolved: bool = False
    replies: List[Dict[str, Any]] = []

class CommentCreate(BaseModel):
    author: str = "Anonymous"
    author_email: Optional[str] = None
    content: str

class CommentReply(BaseModel):
    author: str = "Anonymous"
    content: str

class Signature(BaseModel):
    id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    name: str
    image_base64: str
    created_at: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))

class SignatureCreate(BaseModel):
    name: str
    image_base64: str

class SignaturePlacement(BaseModel):
    signature_id: str
    page: int = 0
    x: float
    y: float
    width: float = 20

class SignatureRequest(BaseModel):
    requester_name: str
    requester_email: str
    signer_email: str
    signer_name: str
    message: Optional[str] = None

class CommentRequest(BaseModel):
    requester_name: str
    requester_email: str
    reviewer_email: str
    reviewer_name: str
    message: Optional[str] = None

class ShareDocumentRequest(BaseModel):
    sender_name: str
    sender_email: str
    recipient_email: str
    recipient_name: str
    message: Optional[str] = None

class PasswordSet(BaseModel):
    password: str

class PasswordVerify(BaseModel):
    password: str

class DocumentCreate(BaseModel):
    document_type: str = "general_document"
    document_subtype: Optional[str] = None
    detected_language: str = "Unknown"
    confidence: float = 0.0
    title: str = "Untitled Document"
    structured_fields: Dict[str, Any] = {}
    formatted_output: str = ""
    tags: List[str] = []
    raw_text: str = ""
    summary: str = ""
    image_thumbnail: Optional[str] = None
    pages_thumbnails: List[str] = []
    pages_count: int = 1
    extracted_dates: List[str] = []
    extracted_amounts: List[str] = []
    extracted_names: List[str] = []
    editor_data: Optional[List[Dict[str, Any]]] = None
    is_edited: bool = False

class DocumentUpdate(BaseModel):
    title: Optional[str] = None
    structured_fields: Optional[Dict[str, Any]] = None
    formatted_output: Optional[str] = None
    summary: Optional[str] = None
    tags: Optional[List[str]] = None
    editor_data: Optional[List[Dict[str, Any]]] = None
    is_edited: Optional[bool] = None
    pages_data: Optional[List[Dict[str, Any]]] = None

class DocumentResponse(DocumentCreate):
    id: str
    created_at: datetime
    size_kb: float = 0.0
    is_locked: bool = False
    comments: List[Dict[str, Any]] = []
    signatures: List[Dict[str, Any]] = []
    signature_requests: List[Dict[str, Any]] = []
    pages_data: List[Dict[str, Any]] = []

# ── Gemini Prompts ────────────────────────────────────────────────────────
SYSTEM_PROMPT = """You are DocScan AI — a world-class multilingual document OCR and classification system.
You support ALL world languages. Extract ALL visible text EXACTLY as written — preserve the original language.
Respond ONLY with valid JSON. No markdown, no code blocks."""

USER_PROMPT = """Analyze this document image and extract ALL information.
Identify the document_type: passport, national_id, drivers_license, invoice, receipt, business_card, contract, bank_statement, medical_record, prescription, handwritten_note, certificate, legal_document, academic_transcript, tax_document, insurance_document, utility_bill, general_document

Return ONLY valid JSON:
{
  "document_type": "...",
  "document_subtype": "...",
  "detected_language": "...",
  "confidence": 0.0-1.0,
  "title": "...",
  "structured_fields": {...},
  "formatted_output": "formatted text with sections",
  "tags": ["tags"],
  "raw_text": "all text",
  "summary": "1-2 sentences",
  "extracted_dates": ["dates"],
  "extracted_amounts": ["amounts"],
  "extracted_names": ["names"]
}"""

AI_ASSISTANT_PROMPT = """You are DocScan AI Assistant — a helpful document analysis assistant.
You help users understand, analyze, summarize, and extract information from their documents.
Be concise but thorough. If you need to reference specific parts of the document, quote them.
Always be helpful and professional."""

def get_multi_prompt(n: int) -> str:
    return f"This document has {n} pages. Analyze ALL pages together.\n\n" + USER_PROMPT

# ── Export Generators ─────────────────────────────────────────────────────
def generate_pdf(doc: dict) -> bytes:
    from fpdf import FPDF
    TYPE_LABELS = {
        'passport': 'PASSPORT', 'national_id': 'NATIONAL ID', 'drivers_license': "DRIVER'S LICENSE",
        'invoice': 'INVOICE', 'receipt': 'RECEIPT', 'business_card': 'BUSINESS CARD',
        'contract': 'CONTRACT', 'bank_statement': 'BANK STATEMENT', 'medical_record': 'MEDICAL RECORD',
        'prescription': 'PRESCRIPTION', 'handwritten_note': 'HANDWRITTEN NOTE', 'certificate': 'CERTIFICATE',
        'legal_document': 'LEGAL DOCUMENT', 'academic_transcript': 'ACADEMIC TRANSCRIPT',
        'tax_document': 'TAX DOCUMENT', 'insurance_document': 'INSURANCE DOCUMENT',
        'utility_bill': 'UTILITY BILL', 'general_document': 'DOCUMENT',
    }
    type_label = TYPE_LABELS.get(doc.get('document_type', ''), 'DOCUMENT')
    pdf = FPDF(orientation='P', unit='mm', format='A4')
    pdf.set_auto_page_break(auto=True, margin=18)
    pdf.add_page()
    pdf.set_fill_color(37, 99, 235)
    pdf.rect(0, 0, 210, 22, 'F')
    pdf.set_font('Helvetica', 'B', 12)
    pdf.set_text_color(255, 255, 255)
    pdf.set_xy(10, 6)
    pdf.cell(100, 10, 'DocScan Pro', ln=False)
    pdf.set_x(110)
    pdf.cell(0, 10, safe_latin(type_label), ln=True, align='R')
    pdf.set_text_color(30, 30, 30)
    pdf.set_xy(15, 30)
    pdf.set_font('Helvetica', 'B', 18)
    pdf.multi_cell(180, 10, safe_latin(doc.get('title', 'Untitled Document')))
    pdf.set_font('Helvetica', '', 9)
    pdf.set_text_color(100, 100, 100)
    pdf.ln(2)
    lang = safe_latin(doc.get('detected_language', 'Unknown'))
    conf = f"{int(doc.get('confidence', 0) * 100)}%"
    pdf.cell(65, 6, f"Language: {lang}")
    pdf.cell(0, 6, f"Confidence: {conf}", ln=True)
    y = pdf.get_y() + 3
    pdf.set_draw_color(37, 99, 235)
    pdf.line(15, y, 195, y)
    pdf.ln(5)
    if doc.get('summary'):
        pdf.set_font('Helvetica', 'I', 10)
        pdf.set_text_color(70, 70, 70)
        pdf.multi_cell(180, 6, safe_latin(doc['summary']))
        pdf.ln(4)
    if doc.get('formatted_output'):
        pdf.set_font('Helvetica', 'B', 12)
        pdf.set_text_color(37, 99, 235)
        pdf.cell(0, 8, 'Extracted Content', ln=True)
        pdf.ln(3)
        pdf.set_font('Courier', '', 8)
        pdf.set_text_color(40, 40, 40)
        for line in doc['formatted_output'].split('\n')[:100]:
            if line.strip():
                pdf.multi_cell(180, 5, safe_latin(line))
            else:
                pdf.ln(2)
    return bytes(pdf.output())

def generate_txt(doc: dict) -> bytes:
    lines = ['DocScan Pro — Document Export', '=' * 60,
             f"Title:      {doc.get('title', '')}",
             f"Type:       {doc.get('document_type', '').replace('_', ' ').title()}",
             f"Language:   {doc.get('detected_language', '')}",
             f"Confidence: {int(doc.get('confidence', 0) * 100)}%", '']
    if doc.get('summary'):
        lines += ['SUMMARY', '-' * 40, doc['summary'], '']
    if doc.get('formatted_output'):
        lines += ['CONTENT', '-' * 40, doc['formatted_output'], '']
    return '\n'.join(lines).encode('utf-8')

def generate_docx(doc: dict) -> bytes:
    from docx import Document as DocxDoc
    d = DocxDoc()
    d.add_heading(doc.get('title', 'Untitled'), 0)
    if doc.get('summary'):
        p = d.add_paragraph(doc['summary'])
        p.runs[0].italic = True
    if doc.get('formatted_output'):
        d.add_heading('Content', 1)
        for line in doc['formatted_output'].split('\n'):
            if line.strip():
                d.add_paragraph(line)
    buf = io.BytesIO()
    d.save(buf)
    return buf.getvalue()

def generate_pptx(doc: dict) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    sl = prs.slides.add_slide(prs.slide_layouts[0])
    sl.shapes.title.text = doc.get('title', 'Document')
    sl.placeholders[1].text = doc.get('summary', '')
    if doc.get('formatted_output'):
        sl2 = prs.slides.add_slide(prs.slide_layouts[1])
        sl2.shapes.title.text = 'Content'
        sl2.placeholders[1].text = doc['formatted_output'][:2000]
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()

def generate_image_export(doc: dict, fmt: str) -> bytes:
    W, H = 900, 1200
    img = Image.new('RGB', (W, H), color=(255, 255, 255))
    draw = ImageDraw.Draw(img)
    draw.rectangle([0, 0, W, 70], fill=(37, 99, 235))
    draw.text((W // 2, 35), 'DocScan Pro', fill=(255, 255, 255), anchor='mm')
    draw.text((W // 2, 100), (doc.get('title', '') or '')[:70], fill=(20, 20, 20), anchor='mt')
    y = 150
    if doc.get('formatted_output'):
        for line in doc['formatted_output'].split('\n')[:30]:
            draw.text((40, y), line[:80], fill=(60, 60, 60))
            y += 20
            if y > H - 100:
                break
    buf = io.BytesIO()
    
    # Handle different image formats
    if fmt.upper() == 'JPEG':
        img.save(buf, format='JPEG', quality=92)
    elif fmt.upper() == 'PNG':
        img.save(buf, format='PNG')
    elif fmt.upper() == 'TIFF':
        img.save(buf, format='TIFF')
    elif fmt.upper() == 'BMP':
        img.save(buf, format='BMP')
    elif fmt.upper() == 'WEBP':
        img.save(buf, format='WEBP', quality=90)
    else:
        img.save(buf, format=fmt.upper())
    return buf.getvalue()

def generate_xlsx(doc: dict) -> bytes:
    """Generate Excel spreadsheet from document"""
    import xlsxwriter
    buf = io.BytesIO()
    wb = xlsxwriter.Workbook(buf, {'in_memory': True})
    
    # Title format
    title_fmt = wb.add_format({'bold': True, 'font_size': 16, 'font_color': '#2563EB'})
    header_fmt = wb.add_format({'bold': True, 'bg_color': '#2563EB', 'font_color': 'white', 'border': 1})
    cell_fmt = wb.add_format({'text_wrap': True, 'valign': 'top', 'border': 1})
    
    # Document Info Sheet
    ws_info = wb.add_worksheet('Document Info')
    ws_info.set_column('A:A', 20)
    ws_info.set_column('B:B', 60)
    
    ws_info.write('A1', doc.get('title', 'Untitled'), title_fmt)
    ws_info.write('A3', 'Property', header_fmt)
    ws_info.write('B3', 'Value', header_fmt)
    
    info_rows = [
        ('Document Type', doc.get('document_type', '').replace('_', ' ').title()),
        ('Language', doc.get('detected_language', 'Unknown')),
        ('Confidence', f"{int(doc.get('confidence', 0) * 100)}%"),
        ('Pages', str(doc.get('pages_count', 1))),
        ('Created', doc.get('created_at', '')),
    ]
    for i, (prop, val) in enumerate(info_rows, start=4):
        ws_info.write(f'A{i}', prop, cell_fmt)
        ws_info.write(f'B{i}', val, cell_fmt)
    
    # Content Sheet
    if doc.get('formatted_output'):
        ws_content = wb.add_worksheet('Content')
        ws_content.set_column('A:A', 100)
        ws_content.write('A1', 'Extracted Content', title_fmt)
        row = 2
        for line in doc['formatted_output'].split('\n'):
            ws_content.write(f'A{row}', line, cell_fmt)
            row += 1
    
    # Summary Sheet
    if doc.get('summary'):
        ws_summary = wb.add_worksheet('Summary')
        ws_summary.set_column('A:A', 80)
        ws_summary.write('A1', 'Document Summary', title_fmt)
        ws_summary.write('A3', doc['summary'], cell_fmt)
    
    wb.close()
    return buf.getvalue()

def generate_html(doc: dict) -> bytes:
    """Generate HTML document"""
    html = f'''<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{doc.get('title', 'Document')}</title>
    <style>
        body {{ font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, sans-serif; max-width: 800px; margin: 0 auto; padding: 20px; background: #f8fafc; }}
        .header {{ background: #2563eb; color: white; padding: 20px; border-radius: 12px; margin-bottom: 20px; }}
        .header h1 {{ margin: 0 0 10px 0; }}
        .meta {{ display: flex; gap: 20px; font-size: 14px; opacity: 0.9; }}
        .section {{ background: white; padding: 20px; border-radius: 12px; margin-bottom: 16px; box-shadow: 0 1px 3px rgba(0,0,0,0.1); }}
        .section h2 {{ color: #2563eb; margin-top: 0; border-bottom: 2px solid #e5e7eb; padding-bottom: 10px; }}
        .content {{ white-space: pre-wrap; font-family: 'SF Mono', Monaco, monospace; font-size: 13px; line-height: 1.6; }}
        .tag {{ display: inline-block; background: #dbeafe; color: #2563eb; padding: 4px 12px; border-radius: 16px; font-size: 12px; margin: 4px; }}
        .footer {{ text-align: center; color: #6b7280; font-size: 12px; margin-top: 20px; }}
    </style>
</head>
<body>
    <div class="header">
        <h1>{doc.get('title', 'Untitled Document')}</h1>
        <div class="meta">
            <span>📄 {doc.get('document_type', 'document').replace('_', ' ').title()}</span>
            <span>🌐 {doc.get('detected_language', 'Unknown')}</span>
            <span>✓ {int(doc.get('confidence', 0) * 100)}% confidence</span>
        </div>
    </div>
'''
    if doc.get('summary'):
        html += f'''    <div class="section">
        <h2>📋 Summary</h2>
        <p>{doc['summary']}</p>
    </div>
'''
    if doc.get('formatted_output'):
        html += f'''    <div class="section">
        <h2>📝 Content</h2>
        <div class="content">{doc['formatted_output']}</div>
    </div>
'''
    if doc.get('tags'):
        tags_html = ''.join([f'<span class="tag">#{tag}</span>' for tag in doc['tags']])
        html += f'''    <div class="section">
        <h2>🏷️ Tags</h2>
        {tags_html}
    </div>
'''
    html += '''    <div class="footer">
        <p>Exported from DocScan Pro</p>
    </div>
</body>
</html>'''
    return html.encode('utf-8')

def generate_json_export(doc: dict) -> bytes:
    """Generate JSON export"""
    # Convert datetime objects to ISO strings
    created_at = doc.get('created_at', '')
    if hasattr(created_at, 'isoformat'):
        created_at = created_at.isoformat()
    elif created_at and not isinstance(created_at, str):
        created_at = str(created_at)
    
    export_data = {
        "document": {
            "title": doc.get('title', ''),
            "type": doc.get('document_type', ''),
            "language": doc.get('detected_language', ''),
            "confidence": doc.get('confidence', 0),
            "pages_count": doc.get('pages_count', 1),
            "created_at": created_at,
        },
        "content": {
            "summary": doc.get('summary', ''),
            "formatted_output": doc.get('formatted_output', ''),
            "raw_text": doc.get('raw_text', ''),
        },
        "metadata": {
            "tags": doc.get('tags', []),
            "is_locked": doc.get('is_locked', False),
            "comments_count": len(doc.get('comments', [])),
            "signatures_count": len(doc.get('signatures', [])),
        },
        "export_info": {
            "exported_at": datetime.now(timezone.utc).isoformat(),
            "source": "DocScan Pro",
            "version": "5.0"
        }
    }
    return json.dumps(export_data, indent=2, ensure_ascii=False).encode('utf-8')

def generate_markdown(doc: dict) -> bytes:
    """Generate Markdown document"""
    md = f'''# {doc.get('title', 'Untitled Document')}

---

**Document Type:** {doc.get('document_type', 'document').replace('_', ' ').title()}  
**Language:** {doc.get('detected_language', 'Unknown')}  
**Confidence:** {int(doc.get('confidence', 0) * 100)}%  
**Pages:** {doc.get('pages_count', 1)}

---

'''
    if doc.get('summary'):
        md += f'''## 📋 Summary

{doc['summary']}

'''
    if doc.get('formatted_output'):
        md += f'''## 📝 Content

```
{doc['formatted_output']}
```

'''
    if doc.get('tags'):
        tags = ' '.join([f'`#{tag}`' for tag in doc['tags']])
        md += f'''## 🏷️ Tags

{tags}

'''
    md += '''---

*Exported from DocScan Pro*
'''
    return md.encode('utf-8')

def generate_epub(doc: dict) -> bytes:
    """Generate EPUB e-book"""
    from ebooklib import epub
    
    book = epub.EpubBook()
    book.set_identifier(doc.get('id', str(uuid.uuid4())))
    book.set_title(doc.get('title', 'Document'))
    book.set_language(doc.get('detected_language', 'en')[:2].lower() or 'en')
    book.add_author('DocScan Pro')
    
    # Create chapter
    c1 = epub.EpubHtml(title=doc.get('title', 'Document'), file_name='content.xhtml', lang='en')
    content_html = f'''<html>
<head><title>{doc.get('title', 'Document')}</title></head>
<body>
<h1>{doc.get('title', 'Document')}</h1>
<p><em>Type: {doc.get('document_type', '').replace('_', ' ').title()}</em></p>
'''
    if doc.get('summary'):
        content_html += f'<h2>Summary</h2><p>{doc["summary"]}</p>'
    if doc.get('formatted_output'):
        content_html += f'<h2>Content</h2><pre>{doc["formatted_output"]}</pre>'
    content_html += '</body></html>'
    c1.content = content_html
    
    book.add_item(c1)
    book.spine = ['nav', c1]
    book.add_item(epub.EpubNcx())
    book.add_item(epub.EpubNav())
    
    buf = io.BytesIO()
    epub.write_epub(buf, book)
    return buf.getvalue()

def generate_mobi(doc: dict) -> bytes:
    """Generate MOBI (Kindle) - returns EPUB as MOBI conversion requires Kindlegen"""
    # Note: True MOBI requires Amazon's Kindlegen tool
    # We return EPUB which Kindle devices can also read
    return generate_epub(doc)

def generate_svg(doc: dict) -> bytes:
    """Generate SVG document visualization"""
    title = doc.get('title', 'Document')[:50]
    doc_type = doc.get('document_type', 'document').replace('_', ' ').title()
    content_lines = (doc.get('formatted_output', '') or '').split('\n')[:20]
    
    svg = f'''<?xml version="1.0" encoding="UTF-8"?>
<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 600 800" width="600" height="800">
  <defs>
    <linearGradient id="headerGrad" x1="0%" y1="0%" x2="100%" y2="0%">
      <stop offset="0%" style="stop-color:#2563eb"/>
      <stop offset="100%" style="stop-color:#3b82f6"/>
    </linearGradient>
  </defs>
  
  <!-- Background -->
  <rect width="600" height="800" fill="#f8fafc"/>
  
  <!-- Header -->
  <rect width="600" height="80" fill="url(#headerGrad)"/>
  <text x="300" y="35" text-anchor="middle" fill="white" font-family="Arial" font-size="18" font-weight="bold">DocScan Pro</text>
  <text x="300" y="60" text-anchor="middle" fill="rgba(255,255,255,0.8)" font-family="Arial" font-size="12">{doc_type}</text>
  
  <!-- Title -->
  <text x="30" y="120" fill="#1e3a5f" font-family="Arial" font-size="20" font-weight="bold">{title}</text>
  
  <!-- Content -->
  <rect x="20" y="140" width="560" height="620" fill="white" stroke="#e5e7eb" rx="8"/>
'''
    y = 170
    for line in content_lines:
        escaped = line[:70].replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
        svg += f'  <text x="35" y="{y}" fill="#374151" font-family="Monaco, monospace" font-size="11">{escaped}</text>\n'
        y += 22
        if y > 720:
            break
    
    svg += '''  
  <!-- Footer -->
  <text x="300" y="780" text-anchor="middle" fill="#9ca3af" font-family="Arial" font-size="10">Exported from DocScan Pro</text>
</svg>'''
    return svg.encode('utf-8')

# ── Routes ────────────────────────────────────────────────────────────────
@api_router.get("/")
async def root():
    return {"message": "DocScan Pro API v5 - Full Featured Document Management"}

@api_router.post("/scan")
async def scan_document(request: ScanRequest):
    api_key = get_api_key()
    images = [strip_b64_prefix(img) for img in request.images]
    n = len(images)
    try:
        chat = LlmChat(api_key=api_key, session_id=str(uuid.uuid4()), system_message=SYSTEM_PROMPT).with_model("gemini", "gemini-3-flash-preview")
        file_contents = [ImageContent(image_base64=img) for img in images]
        prompt = get_multi_prompt(n) if n > 1 else USER_PROMPT
        response = await chat.send_message(UserMessage(text=prompt, file_contents=file_contents))
        try:
            result = json.loads(response)
        except json.JSONDecodeError:
            m = re.search(r'\{[\s\S]*\}', response)
            if m:
                result = json.loads(m.group())
            else:
                raise HTTPException(500, "AI returned non-JSON response")
        result['pages_count'] = n
        return result
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Scan error: {e}")
        raise HTTPException(500, f"AI analysis failed: {str(e)}")

# ── Image Processing Routes ────────────────────────────────────────────────
@api_router.post("/process-image")
async def process_image(request: ImageProcessRequest):
    """Process image with various operations"""
    try:
        img = b64_to_cv2(request.image)
        
        for op in request.operations:
            if op == 'straighten':
                img = straighten_document(img)
            elif op == 'enhance':
                img = enhance_document(img)
        
        result_b64 = cv2_to_b64(img)
        
        edges = []
        if 'detect_edges' in request.operations:
            edges = detect_document_edges(img)
        
        return {
            "image": result_b64,
            "edges": edges
        }
    except Exception as e:
        logger.error(f"Image processing error: {e}")
        raise HTTPException(500, f"Image processing failed: {str(e)}")

@api_router.post("/perspective-transform")
async def perspective_transform(image: str, points: List[List[float]]):
    """Apply perspective transform to flatten document"""
    try:
        img = b64_to_cv2(image)
        transformed = apply_perspective_transform(img, points)
        return {"image": cv2_to_b64(transformed)}
    except Exception as e:
        raise HTTPException(500, f"Transform failed: {str(e)}")

@api_router.post("/add-watermark")
async def add_watermark_route(request: WatermarkRequest):
    """Add watermark to image"""
    try:
        img = b64_to_pil(request.image)
        result = add_watermark(img, request.text, request.opacity)
        return {"image": pil_to_b64(result.convert('RGB'))}
    except Exception as e:
        raise HTTPException(500, f"Watermark failed: {str(e)}")

@api_router.post("/add-blur")
async def add_blur_route(request: BlurRequest):
    """Blur a region of the image"""
    try:
        img = b64_to_pil(request.image)
        result = add_blur_region(img, request.x, request.y, request.width, request.height, request.intensity)
        return {"image": pil_to_b64(result)}
    except Exception as e:
        raise HTTPException(500, f"Blur failed: {str(e)}")

# ── AI Assistant Route ─────────────────────────────────────────────────────
@api_router.post("/ai-assistant")
async def ai_assistant(request: AIAssistantRequest):
    """AI assistant for document questions"""
    api_key = get_api_key()
    
    context = request.context or ""
    
    # If document_id provided, fetch document context
    if request.document_id:
        doc = await db.documents.find_one({"id": request.document_id})
        if doc:
            context = f"Document: {doc.get('title', 'Untitled')}\nType: {doc.get('document_type', 'Unknown')}\nContent: {doc.get('formatted_output', '')[:2000]}\n\n"
    
    try:
        chat = LlmChat(
            api_key=api_key,
            session_id=str(uuid.uuid4()),
            system_message=AI_ASSISTANT_PROMPT
        ).with_model("gemini", "gemini-3-flash-preview")
        
        full_message = f"{context}User question: {request.message}"
        response = await chat.send_message(UserMessage(text=full_message))
        
        return {"response": response}
    except Exception as e:
        logger.error(f"AI Assistant error: {e}")
        raise HTTPException(500, f"AI Assistant failed: {str(e)}")

# ── OCR Text Recognition ───────────────────────────────────────────────────
@api_router.post("/recognize-text")
async def recognize_text(image: str = "", region: Optional[Dict[str, int]] = None):
    """Recognize text in image or specific region using AI"""
    api_key = get_api_key()
    
    try:
        img_data = strip_b64_prefix(image)
        
        # If region specified, crop the image
        if region:
            pil_img = b64_to_pil(image)
            cropped = pil_img.crop((region['x'], region['y'], region['x'] + region['width'], region['y'] + region['height']))
            buf = io.BytesIO()
            cropped.save(buf, format='JPEG', quality=90)
            img_data = base64.b64encode(buf.getvalue()).decode()
        
        chat = LlmChat(
            api_key=api_key,
            session_id=str(uuid.uuid4()),
            system_message="Extract all text from this image exactly as shown. Return only the text, no formatting."
        ).with_model("gemini", "gemini-3-flash-preview")
        
        response = await chat.send_message(UserMessage(
            text="Extract all text from this image:",
            file_contents=[ImageContent(image_base64=img_data)]
        ))
        
        return {"text": response}
    except Exception as e:
        raise HTTPException(500, f"Text recognition failed: {str(e)}")

# ── Measurement Modes ──────────────────────────────────────────────────────
@api_router.post("/measure")
async def measure(request: MeasurementRequest):
    """Count objects or measure areas in document"""
    api_key = get_api_key()
    
    try:
        if request.mode == 'count':
            # Use AI to count items
            chat = LlmChat(
                api_key=api_key,
                session_id=str(uuid.uuid4()),
                system_message="You are a counting assistant. Count the specified items in images."
            ).with_model("gemini", "gemini-3-flash-preview")
            
            response = await chat.send_message(UserMessage(
                text="Count all distinct items/objects in this image. Return JSON: {\"count\": number, \"items\": [list of items]}",
                file_contents=[ImageContent(image_base64=strip_b64_prefix(request.image))]
            ))
            
            try:
                result = json.loads(response)
            except:
                result = {"count": 0, "raw_response": response}
            
            return result
            
        elif request.mode == 'area':
            # Calculate area from points
            if request.points and len(request.points) >= 3:
                # Shoelace formula for polygon area
                n = len(request.points)
                area = 0.0
                for i in range(n):
                    j = (i + 1) % n
                    area += request.points[i]['x'] * request.points[j]['y']
                    area -= request.points[j]['x'] * request.points[i]['y']
                area = abs(area) / 2.0
                return {"area_pixels": area, "points": request.points}
            
            return {"error": "Need at least 3 points for area calculation"}
        
        return {"error": f"Unknown mode: {request.mode}"}
    except Exception as e:
        raise HTTPException(500, f"Measurement failed: {str(e)}")

# ── Cloud Storage Routes ───────────────────────────────────────────────────
@api_router.post("/cloud/connect")
async def cloud_connect(provider: str):
    """Get OAuth URL for cloud provider"""
    # This would normally return OAuth URLs for each provider
    providers = {
        'google_drive': {'name': 'Google Drive', 'icon': 'logo-google'},
        'dropbox': {'name': 'Dropbox', 'icon': 'cloud-outline'},
        'onedrive': {'name': 'OneDrive', 'icon': 'logo-microsoft'},
        'box': {'name': 'Box', 'icon': 'cube-outline'},
        'icloud': {'name': 'iCloud', 'icon': 'logo-apple'},
    }
    
    if provider not in providers:
        raise HTTPException(400, f"Unknown provider: {provider}")
    
    # In production, return actual OAuth URL
    return {
        "provider": provider,
        "auth_url": f"https://oauth.{provider}.com/authorize",
        "status": "oauth_required"
    }

@api_router.get("/cloud/providers")
async def list_cloud_providers():
    """List available cloud storage providers"""
    return {
        "providers": [
            {"id": "google_drive", "name": "Google Drive", "icon": "logo-google", "connected": False},
            {"id": "dropbox", "name": "Dropbox", "icon": "cloud-outline", "connected": False},
            {"id": "onedrive", "name": "OneDrive", "icon": "logo-microsoft", "connected": False},
            {"id": "box", "name": "Box", "icon": "cube-outline", "connected": False},
            {"id": "icloud", "name": "iCloud", "icon": "logo-apple", "connected": False},
        ]
    }

# ── Document CRUD ──────────────────────────────────────────────────────────
@api_router.post("/documents", response_model=DocumentResponse)
async def create_document(doc: DocumentCreate):
    doc_id = str(uuid.uuid4())
    now = datetime.now(timezone.utc)
    size_kb = round(len(doc.image_thumbnail or '') * 3 / 4 / 1024, 1)
    doc_dict = {
        "id": doc_id, "created_at": now, "size_kb": size_kb,
        "is_locked": False, "password_hash": None,
        "comments": [], "signatures": [], "signature_requests": [],
        "pages_data": [],
        **doc.dict()
    }
    await db.documents.insert_one(doc_dict)
    return DocumentResponse(**{k: v for k, v in doc_dict.items() if k != 'password_hash'})

@api_router.get("/documents", response_model=List[DocumentResponse])
async def list_documents():
    docs = await db.documents.find({}, {"_id": 0, "password_hash": 0}).sort("created_at", -1).to_list(500)
    return [DocumentResponse(**d) for d in docs]

@api_router.get("/documents/{doc_id}", response_model=DocumentResponse)
async def get_document(doc_id: str):
    doc = await db.documents.find_one({"id": doc_id}, {"_id": 0, "password_hash": 0})
    if not doc:
        raise HTTPException(404, "Document not found")
    return DocumentResponse(**doc)

@api_router.put("/documents/{doc_id}", response_model=DocumentResponse)
async def update_document(doc_id: str, updates: DocumentUpdate):
    payload = {k: v for k, v in updates.dict().items() if v is not None}
    if not payload:
        raise HTTPException(400, "No fields to update")
    result = await db.documents.update_one({"id": doc_id}, {"$set": payload})
    if result.matched_count == 0:
        raise HTTPException(404, "Document not found")
    updated = await db.documents.find_one({"id": doc_id}, {"_id": 0, "password_hash": 0})
    return DocumentResponse(**updated)

@api_router.delete("/documents/{doc_id}")
async def delete_document(doc_id: str):
    result = await db.documents.delete_one({"id": doc_id})
    if result.deleted_count == 0:
        raise HTTPException(404, "Document not found")
    return {"message": "Deleted"}

# ── Password Protection ────────────────────────────────────────────────────
@api_router.post("/documents/{doc_id}/password")
async def set_document_password(doc_id: str, data: PasswordSet):
    doc = await db.documents.find_one({"id": doc_id})
    if not doc:
        raise HTTPException(404, "Document not found")
    password_hash = hash_password(data.password)
    await db.documents.update_one({"id": doc_id}, {"$set": {"password_hash": password_hash, "is_locked": True}})
    return {"message": "Password set", "is_locked": True}

@api_router.post("/documents/{doc_id}/verify-password")
async def verify_document_password(doc_id: str, data: PasswordVerify):
    doc = await db.documents.find_one({"id": doc_id})
    if not doc:
        raise HTTPException(404, "Document not found")
    if not doc.get('is_locked') or not doc.get('password_hash'):
        return {"verified": True}
    if verify_password(data.password, doc['password_hash']):
        return {"verified": True}
    raise HTTPException(403, "Incorrect password")

@api_router.delete("/documents/{doc_id}/password")
async def remove_document_password(doc_id: str, data: PasswordVerify):
    doc = await db.documents.find_one({"id": doc_id})
    if not doc:
        raise HTTPException(404, "Document not found")
    if doc.get('password_hash') and not verify_password(data.password, doc['password_hash']):
        raise HTTPException(403, "Incorrect password")
    await db.documents.update_one({"id": doc_id}, {"$set": {"password_hash": None, "is_locked": False}})
    return {"message": "Password removed", "is_locked": False}

# ── Comments ───────────────────────────────────────────────────────────────
@api_router.post("/documents/{doc_id}/comments")
async def add_comment(doc_id: str, comment: CommentCreate):
    doc = await db.documents.find_one({"id": doc_id})
    if not doc:
        raise HTTPException(404, "Document not found")
    new_comment = {
        "id": str(uuid.uuid4()),
        "author": comment.author,
        "author_email": comment.author_email,
        "content": comment.content,
        "created_at": datetime.now(timezone.utc).isoformat(),
        "resolved": False,
        "replies": []
    }
    await db.documents.update_one({"id": doc_id}, {"$push": {"comments": new_comment}})
    return {"message": "Comment added", "comment": new_comment}

@api_router.post("/documents/{doc_id}/comments/{comment_id}/reply")
async def reply_to_comment(doc_id: str, comment_id: str, reply: CommentReply):
    new_reply = {
        "id": str(uuid.uuid4()),
        "author": reply.author,
        "content": reply.content,
        "created_at": datetime.now(timezone.utc).isoformat()
    }
    result = await db.documents.update_one(
        {"id": doc_id, "comments.id": comment_id},
        {"$push": {"comments.$.replies": new_reply}}
    )
    if result.modified_count == 0:
        raise HTTPException(404, "Comment not found")
    return {"message": "Reply added", "reply": new_reply}

@api_router.put("/documents/{doc_id}/comments/{comment_id}/resolve")
async def resolve_comment(doc_id: str, comment_id: str):
    result = await db.documents.update_one(
        {"id": doc_id, "comments.id": comment_id},
        {"$set": {"comments.$.resolved": True}}
    )
    if result.modified_count == 0:
        raise HTTPException(404, "Comment not found")
    return {"message": "Comment resolved"}

@api_router.delete("/documents/{doc_id}/comments/{comment_id}")
async def delete_comment(doc_id: str, comment_id: str):
    result = await db.documents.update_one({"id": doc_id}, {"$pull": {"comments": {"id": comment_id}}})
    if result.modified_count == 0:
        raise HTTPException(404, "Comment not found")
    return {"message": "Comment deleted"}

@api_router.post("/documents/{doc_id}/request-comment")
async def request_comment(doc_id: str, request: CommentRequest, background_tasks: BackgroundTasks):
    doc = await db.documents.find_one({"id": doc_id})
    if not doc:
        raise HTTPException(404, "Document not found")
    
    # Send email via Resend
    background_tasks.add_task(
        send_email,
        request.reviewer_email,
        'comment_request',
        {
            'reviewer_name': request.reviewer_name,
            'requester_name': request.requester_name,
            'requester_email': request.requester_email,
            'document_title': doc.get('title', 'Untitled'),
            'message': request.message,
            'action_url': f"https://docscan.app/document/{doc_id}"
        }
    )
    
    return {"message": f"Comment request sent to {request.reviewer_email}"}

# ── Signatures ─────────────────────────────────────────────────────────────
@api_router.post("/signatures")
async def create_signature(signature: SignatureCreate):
    sig_id = str(uuid.uuid4())
    sig_data = {
        "id": sig_id,
        "name": signature.name,
        "image_base64": signature.image_base64,
        "created_at": datetime.now(timezone.utc).isoformat()
    }
    await db.signatures.insert_one(sig_data)
    return {"message": "Signature saved", "signature": {k: v for k, v in sig_data.items() if k != '_id'}}

@api_router.get("/signatures")
async def list_signatures():
    sigs = await db.signatures.find({}, {"_id": 0}).sort("created_at", -1).to_list(100)
    return sigs

@api_router.delete("/signatures/{sig_id}")
async def delete_signature(sig_id: str):
    result = await db.signatures.delete_one({"id": sig_id})
    if result.deleted_count == 0:
        raise HTTPException(404, "Signature not found")
    return {"message": "Signature deleted"}

@api_router.post("/documents/{doc_id}/signatures")
async def add_signature_to_document(doc_id: str, placement: SignaturePlacement):
    doc = await db.documents.find_one({"id": doc_id})
    if not doc:
        raise HTTPException(404, "Document not found")
    sig = await db.signatures.find_one({"id": placement.signature_id})
    if not sig:
        raise HTTPException(404, "Signature not found")
    sig_placement = {
        "id": str(uuid.uuid4()),
        "signature_id": placement.signature_id,
        "signature_name": sig.get('name'),
        "signature_image": sig.get('image_base64'),
        "page": placement.page,
        "x": placement.x,
        "y": placement.y,
        "width": placement.width,
        "placed_at": datetime.now(timezone.utc).isoformat()
    }
    await db.documents.update_one({"id": doc_id}, {"$push": {"signatures": sig_placement}})
    return {"message": "Signature added", "placement": sig_placement}

@api_router.delete("/documents/{doc_id}/signatures/{placement_id}")
async def remove_signature_from_document(doc_id: str, placement_id: str):
    result = await db.documents.update_one({"id": doc_id}, {"$pull": {"signatures": {"id": placement_id}}})
    if result.modified_count == 0:
        raise HTTPException(404, "Signature placement not found")
    return {"message": "Signature removed"}

@api_router.post("/documents/{doc_id}/request-signature")
async def request_signature(doc_id: str, request: SignatureRequest, background_tasks: BackgroundTasks):
    doc = await db.documents.find_one({"id": doc_id})
    if not doc:
        raise HTTPException(404, "Document not found")
    
    sig_request = {
        "id": str(uuid.uuid4()),
        "requester_name": request.requester_name,
        "requester_email": request.requester_email,
        "signer_email": request.signer_email,
        "signer_name": request.signer_name,
        "message": request.message,
        "status": "pending",
        "created_at": datetime.now(timezone.utc).isoformat()
    }
    
    await db.documents.update_one({"id": doc_id}, {"$push": {"signature_requests": sig_request}})
    
    # Send email via Resend
    background_tasks.add_task(
        send_email,
        request.signer_email,
        'signature_request',
        {
            'signer_name': request.signer_name,
            'requester_name': request.requester_name,
            'requester_email': request.requester_email,
            'document_title': doc.get('title', 'Untitled'),
            'message': request.message,
            'action_url': f"https://docscan.app/sign/{doc_id}"
        }
    )
    
    return {"message": "Signature request sent", "request": sig_request}

# ── Share Document ─────────────────────────────────────────────────────────
@api_router.post("/documents/{doc_id}/share")
async def share_document(doc_id: str, request: ShareDocumentRequest, background_tasks: BackgroundTasks):
    doc = await db.documents.find_one({"id": doc_id})
    if not doc:
        raise HTTPException(404, "Document not found")
    
    # Send email via Resend
    background_tasks.add_task(
        send_email,
        request.recipient_email,
        'document_shared',
        {
            'recipient_name': request.recipient_name,
            'sender_name': request.sender_name,
            'sender_email': request.sender_email,
            'document_title': doc.get('title', 'Untitled'),
            'message': request.message,
            'action_url': f"https://docscan.app/document/{doc_id}"
        }
    )
    
    return {"message": f"Document shared with {request.recipient_email}"}

# ── Export ─────────────────────────────────────────────────────────────────
@api_router.post("/documents/{doc_id}/export")
async def export_document(doc_id: str, format: str = Query("pdf")):
    doc = await db.documents.find_one({"id": doc_id}, {"_id": 0})
    if not doc:
        raise HTTPException(404, "Document not found")
    fmt = format.lower().lstrip('.')
    try:
        # Document formats
        if fmt == "pdf":
            data, mime, ext = generate_pdf(doc), "application/pdf", "pdf"
        elif fmt in ("docx", "doc"):
            data = generate_docx(doc)
            mime, ext = "application/vnd.openxmlformats-officedocument.wordprocessingml.document", "docx"
        elif fmt == "pptx":
            data = generate_pptx(doc)
            mime, ext = "application/vnd.openxmlformats-officedocument.presentationml.presentation", "pptx"
        elif fmt in ("xlsx", "xls", "excel"):
            data = generate_xlsx(doc)
            mime, ext = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", "xlsx"
        elif fmt == "txt":
            data, mime, ext = generate_txt(doc), "text/plain", "txt"
        elif fmt == "html":
            data, mime, ext = generate_html(doc), "text/html", "html"
        elif fmt == "json":
            data, mime, ext = generate_json_export(doc), "application/json", "json"
        elif fmt in ("md", "markdown"):
            data, mime, ext = generate_markdown(doc), "text/markdown", "md"
        
        # Image formats
        elif fmt in ("jpeg", "jpg"):
            data, mime, ext = generate_image_export(doc, "JPEG"), "image/jpeg", "jpg"
        elif fmt == "png":
            data, mime, ext = generate_image_export(doc, "PNG"), "image/png", "png"
        elif fmt == "tiff":
            data, mime, ext = generate_image_export(doc, "TIFF"), "image/tiff", "tiff"
        elif fmt == "bmp":
            data, mime, ext = generate_image_export(doc, "BMP"), "image/bmp", "bmp"
        elif fmt == "webp":
            data, mime, ext = generate_image_export(doc, "WEBP"), "image/webp", "webp"
        elif fmt == "svg":
            data, mime, ext = generate_svg(doc), "image/svg+xml", "svg"
        
        # E-book formats
        elif fmt == "epub":
            data, mime, ext = generate_epub(doc), "application/epub+zip", "epub"
        elif fmt == "mobi":
            data, mime, ext = generate_mobi(doc), "application/x-mobipocket-ebook", "mobi"
        
        else:
            raise HTTPException(400, f"Unsupported format: {fmt}. Supported: pdf, docx, pptx, xlsx, txt, html, json, md, jpg, png, tiff, bmp, webp, svg, epub, mobi")
        
        safe_title = re.sub(r'[^\w\s-]', '', doc.get('title', 'document')).strip()[:40]
        filename = f"{safe_title.replace(' ', '_') or 'document'}.{ext}"
        return {"base64": base64.b64encode(data).decode(), "mime_type": mime, "filename": filename}
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Export error: {e}")
        raise HTTPException(500, f"Export failed: {str(e)}")

@api_router.get("/stats")
async def get_stats():
    total = await db.documents.count_documents({})
    locked = await db.documents.count_documents({"is_locked": True})
    recent = await db.documents.find_one({}, {"_id": 0, "created_at": 1}, sort=[("created_at", -1)])
    last_scan = "Never"
    if recent and recent.get("created_at"):
        created = recent["created_at"]
        if isinstance(created, str):
            created = datetime.fromisoformat(created.replace('Z', '+00:00'))
        if created.tzinfo is None:
            created = created.replace(tzinfo=timezone.utc)
        delta = datetime.now(timezone.utc) - created
        hours = int(delta.total_seconds() / 3600)
        last_scan = "Just now" if hours < 1 else (f"{hours}h ago" if hours < 24 else f"{hours // 24}d ago")
    storage_kb = total * 150
    storage_str = f"{storage_kb / 1024:.1f} MB" if storage_kb >= 1024 else f"{storage_kb} KB"
    return {"total_scans": total, "locked_documents": locked, "storage_used": storage_str, "last_scan": last_scan}


# ── Math Solver ────────────────────────────────────────────────────────────
class MathSolveRequest(BaseModel):
    image_base64: Optional[str] = None
    equation: Optional[str] = None

@api_router.post("/math/solve")
async def solve_math_problem(request: MathSolveRequest):
    """
    Solve math problems from image or text input using Gemini AI.
    Supports both image-based (photo of math problem) and text-based equations.
    """
    if not request.image_base64 and not request.equation:
        raise HTTPException(400, "Please provide either an image or an equation")
    
    try:
        api_key = get_api_key()
        chat = LlmChat(
            api_key=api_key,
            session_id=f"math-solver-{uuid.uuid4()}",
            system_message="""You are an expert math tutor and problem solver. Your task is to:
1. Identify the mathematical problem (from image or text)
2. Solve it step-by-step
3. Provide a clear, educational explanation

Format your response as:
**Problem:** [State the problem clearly]

**Solution:**
Step 1: [First step with explanation]
Step 2: [Second step with explanation]
...

**Answer:** [Final answer]

**Explanation:** [Brief explanation of the concept or method used]

Be thorough but concise. Use proper mathematical notation where possible."""
        ).with_model("gemini", "gemini-2.5-flash")
        
        if request.image_base64:
            # Image-based math solving
            clean_base64 = strip_b64_prefix(request.image_base64)
            image_content = ImageContent(image_base64=clean_base64)
            
            user_message = UserMessage(
                text="Please analyze this math problem image and solve it step-by-step. Show all your work and explain each step clearly.",
                image_contents=[image_content]
            )
        else:
            # Text-based equation solving
            user_message = UserMessage(
                text=f"Please solve this math problem step-by-step: {request.equation}\n\nShow all your work and explain each step clearly."
            )
        
        response = await chat.send_message(user_message)
        
        return {
            "success": True,
            "solution": response,
            "input_type": "image" if request.image_base64 else "text",
            "original_equation": request.equation if request.equation else "Extracted from image"
        }
        
    except Exception as e:
        logger.error(f"Math solver error: {e}")
        raise HTTPException(500, f"Failed to solve math problem: {str(e)}")


# ── AI Document Categorization ────────────────────────────────────────────────
class CategorizeRequest(BaseModel):
    text: str
    image_base64: Optional[str] = None

DOCUMENT_CATEGORIES = [
    "Invoice", "Receipt", "Contract", "Letter", "Resume", "ID Document",
    "Bank Statement", "Tax Form", "Medical Record", "Legal Document",
    "Certificate", "Report", "Meeting Notes", "Handwritten Notes",
    "Business Card", "Form", "Other"
]

@api_router.post("/categorize")
async def categorize_document(request: CategorizeRequest):
    """AI-powered document categorization"""
    try:
        api_key = get_api_key()
        chat = LlmChat(
            api_key=api_key,
            session_id=f"categorize-{uuid.uuid4()}",
            system_message=f"""You are a document classification expert. Analyze the document content and classify it into one of these categories:
{', '.join(DOCUMENT_CATEGORIES)}

Also extract:
- Key information (dates, amounts, names, etc.)
- Suggested tags
- Confidence level (high/medium/low)

Return JSON format:
{{
    "category": "Category Name",
    "confidence": "high|medium|low",
    "key_info": {{"field": "value"}},
    "suggested_tags": ["tag1", "tag2"],
    "summary": "Brief document summary"
}}"""
        ).with_model("gemini", "gemini-2.5-flash")
        
        if request.image_base64:
            clean_base64 = strip_b64_prefix(request.image_base64)
            image_content = ImageContent(image_base64=clean_base64)
            user_message = UserMessage(
                text="Classify this document and extract key information.",
                image_contents=[image_content]
            )
        else:
            user_message = UserMessage(
                text=f"Classify this document and extract key information:\n\n{request.text[:2000]}"
            )
        
        response = await chat.send_message(user_message)
        
        # Try to parse as JSON
        try:
            import json
            # Extract JSON from response
            json_match = re.search(r'\{[\s\S]*\}', response)
            if json_match:
                result = json.loads(json_match.group())
            else:
                result = {
                    "category": "Other",
                    "confidence": "low",
                    "key_info": {},
                    "suggested_tags": [],
                    "summary": response[:200]
                }
        except:
            result = {
                "category": "Other",
                "confidence": "low",
                "key_info": {},
                "suggested_tags": [],
                "summary": response[:200]
            }
        
        return {"success": True, **result}
        
    except Exception as e:
        logger.error(f"Categorization error: {e}")
        raise HTTPException(500, f"Failed to categorize document: {str(e)}")


# ── Batch Scanning Endpoint ────────────────────────────────────────────────
class BatchScanRequest(BaseModel):
    images: List[str]  # List of base64 images
    title_prefix: str = "Batch Scan"
    auto_categorize: bool = True

@api_router.post("/batch-scan")
async def batch_scan_documents(request: BatchScanRequest):
    """Process multiple scanned images as a batch"""
    if not request.images:
        raise HTTPException(400, "No images provided")
    
    results = []
    for i, image_b64 in enumerate(request.images):
        try:
            # Process each image
            clean_b64 = strip_b64_prefix(image_b64)
            
            # OCR extraction
            api_key = get_api_key()
            chat = LlmChat(
                api_key=api_key,
                session_id=f"batch-{uuid.uuid4()}",
                system_message="Extract all text from this document image accurately."
            ).with_model("gemini", "gemini-2.5-flash")
            
            image_content = ImageContent(image_base64=clean_b64)
            user_message = UserMessage(
                text="Extract all text from this document.",
                image_contents=[image_content]
            )
            
            extracted_text = await chat.send_message(user_message)
            
            # Create document
            doc_id = str(uuid.uuid4())
            doc = {
                "_id": doc_id,
                "title": f"{request.title_prefix} - Page {i + 1}",
                "raw_text": extracted_text,
                "formatted_output": extracted_text,
                "pages": [{"page_number": 1, "image_base64": clean_b64, "text": extracted_text}],
                "tags": ["batch-scan"],
                "scannedAt": datetime.now(timezone.utc),
                "batch_id": request.title_prefix.replace(" ", "_").lower() + f"_{int(datetime.now().timestamp())}",
            }
            
            await db.documents.insert_one(doc)
            
            results.append({
                "page": i + 1,
                "document_id": doc_id,
                "title": doc["title"],
                "text_preview": extracted_text[:200] + "..." if len(extracted_text) > 200 else extracted_text,
                "success": True
            })
            
        except Exception as e:
            results.append({
                "page": i + 1,
                "success": False,
                "error": str(e)
            })
    
    return {
        "total_pages": len(request.images),
        "successful": sum(1 for r in results if r.get("success")),
        "failed": sum(1 for r in results if not r.get("success")),
        "results": results
    }


# ── Advanced Search/Filter Endpoint ────────────────────────────────────────────────
class AdvancedSearchRequest(BaseModel):
    query: Optional[str] = None
    tags: Optional[List[str]] = None
    category: Optional[str] = None
    date_from: Optional[str] = None
    date_to: Optional[str] = None
    has_password: Optional[bool] = None
    has_signature: Optional[bool] = None
    sort_by: str = "scannedAt"
    sort_order: str = "desc"
    page: int = 1
    limit: int = 20

@api_router.post("/documents/search")
async def advanced_search(request: AdvancedSearchRequest):
    """Advanced document search with filters"""
    query = {}
    
    # Text search
    if request.query:
        query["$or"] = [
            {"title": {"$regex": request.query, "$options": "i"}},
            {"raw_text": {"$regex": request.query, "$options": "i"}},
            {"summary": {"$regex": request.query, "$options": "i"}},
        ]
    
    # Tag filter
    if request.tags:
        query["tags"] = {"$all": request.tags}
    
    # Category filter
    if request.category:
        query["category"] = request.category
    
    # Date range filter
    if request.date_from or request.date_to:
        date_query = {}
        if request.date_from:
            date_query["$gte"] = datetime.fromisoformat(request.date_from.replace("Z", "+00:00"))
        if request.date_to:
            date_query["$lte"] = datetime.fromisoformat(request.date_to.replace("Z", "+00:00"))
        query["scannedAt"] = date_query
    
    # Password filter
    if request.has_password is not None:
        if request.has_password:
            query["password_hash"] = {"$exists": True, "$ne": None}
        else:
            query["$or"] = query.get("$or", []) + [
                {"password_hash": {"$exists": False}},
                {"password_hash": None}
            ]
    
    # Signature filter
    if request.has_signature is not None:
        if request.has_signature:
            query["signatures"] = {"$exists": True, "$ne": []}
        else:
            query["$or"] = query.get("$or", []) + [
                {"signatures": {"$exists": False}},
                {"signatures": []}
            ]
    
    # Sort
    sort_direction = -1 if request.sort_order == "desc" else 1
    
    # Execute query with pagination
    skip = (request.page - 1) * request.limit
    cursor = db.documents.find(query, {"_id": 1, "title": 1, "tags": 1, "scannedAt": 1, "category": 1})
    cursor = cursor.sort(request.sort_by, sort_direction).skip(skip).limit(request.limit)
    
    docs = await cursor.to_list(length=request.limit)
    total = await db.documents.count_documents(query)
    
    # Format results
    results = []
    for doc in docs:
        results.append({
            "id": str(doc["_id"]),
            "title": doc.get("title", "Untitled"),
            "tags": doc.get("tags", []),
            "category": doc.get("category"),
            "scannedAt": doc.get("scannedAt").isoformat() if doc.get("scannedAt") else None,
        })
    
    return {
        "results": results,
        "total": total,
        "page": request.page,
        "limit": request.limit,
        "total_pages": (total + request.limit - 1) // request.limit
    }

app.include_router(api_router)
app.include_router(auth_router, prefix="/api")
app.include_router(subscription_router, prefix="/api")
app.add_middleware(CORSMiddleware, allow_credentials=True, allow_origins=["*"], allow_methods=["*"], allow_headers=["*"])

@app.on_event("shutdown")
async def shutdown_db_client():
    client.close()
